from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import glob
import h5py
import numpy as np
import healpy as hp
import os

def add_image_full_slide(slide, image_path, slide_width, slide_height, title=None):
    with Image.open(image_path) as img:
        img_width, img_height = img.size

    # reserve top strip for title if needed
    top_margin = Inches(0.4) if title else 0
    usable_height = slide_height - top_margin

    img_aspect = img_width / img_height
    slide_aspect = slide_width / usable_height

    if img_aspect > slide_aspect:
        width = slide_width
        height = width / img_aspect
    else:
        height = usable_height
        width = height * img_aspect

    left = (slide_width - width) / 2
    top = top_margin + (usable_height - height) / 2

    slide.shapes.add_picture(image_path, left, top, width=width, height=height)

    if title:
        txBox = slide.shapes.add_textbox(Inches(0), Inches(0), slide_width, top_margin)
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.runs[0].font.size = Pt(16)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = RGBColor(255, 255, 255)


def pngs_to_pptx(slides_spec, output_pptx="output_slides.pptx", area_slide=True):
    """
    slides_spec: list of (image_path, title_or_None)
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- add area slide FIRST ---
    if area_slide:
        add_text_slide(prs, area_slide["title"], area_slide["content"])
    
    for image_path, title in slides_spec:
        if not os.path.exists(image_path):
            print(f"Skipping missing image: {image_path}")
            continue
        slide = prs.slides.add_slide(blank_layout)
        add_image_full_slide(slide, image_path, prs.slide_width, prs.slide_height, title=title)

    prs.save(output_pptx)
    print(f"Saved {output_pptx}")

def read_area_info(h5_path):
    with h5py.File(h5_path, "r") as f:
        area = f.attrs["area_deg2"]
        nside = f.attrs["nside"]
    return area, nside

def nside_to_arcmin(nside):
    return hp.nside2resol(nside, arcmin=True)

def add_text_slide(prs, title, content):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width, Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), prs.slide_width, prs.slide_height)
    tf = content_box.text_frame
    tf.text = content

    return slide

def make_slides(plot_dir):

    # --- read area stats ---
    withcuts_file = plot_dir + "HealpixStats_nside512_withcuts.h5"
    nocuts_file   = plot_dir + "HealpixStats_nside512_nocuts.h5"

    area_withcuts, nside = read_area_info(withcuts_file)
    area_nocuts, _       = read_area_info(nocuts_file)

    resol = nside_to_arcmin(nside)

    area_text = (
        f"Area of the occupied pixels at nside={nside} "
        f"({resol:.2f} arcmin pixel size)\n\n"
        f"With cuts (must have good Z_B)     {area_withcuts:.2f} deg²\n"
        f"Without cuts (Z_B can be missing) {area_nocuts:.2f} deg²"
    )

    area_slide = {
        "title": "Area",
        "content": area_text
    }

    # --- histograms ---
    hist_nocuts_photoz  = plot_dir + "hist_nocuts_hist_photoz.png"
    hist_withcuts_photoz = plot_dir + "hist_withcuts_hist_photoz.png"

    other_hist_withcuts = sorted([
        f for f in glob.glob(plot_dir + "hist_withcuts_hist_*.png")
        if "photoz" not in f
    ])

    # --- healpix maps ---
    hp_nocuts_ZB    = plot_dir + "hp_mean_nocuts_Z_B.png"
    hp_withcuts_all = sorted([
        f for f in glob.glob(plot_dir + "hp_mean_withcuts*.png")
        if "zoom" not in f
    ])

    hp_withcuts_zb = plot_dir + "hp_mean_withcuts_Z_B.png"
    other_hp_withcuts = [f for f in hp_withcuts_all if "Z_B" not in f]

    slides_spec = (
        [(hist_nocuts_photoz,   "no cuts")]
      + [(hist_withcuts_photoz, "with cuts")]
      + [(f, "with cuts") for f in other_hist_withcuts]
      + [(hp_nocuts_ZB,         "no cuts")]
      + [(hp_withcuts_zb,       "with cuts")]
      + [(f, "with cuts") for f in other_hp_withcuts]
    )

    pngs_to_pptx(slides_spec, plot_dir + "plots.pptx", area_slide=area_slide)